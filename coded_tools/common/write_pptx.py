"""
coded_tools/common/write_pptx.py
â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
WritePptx â€” create a Cognizant-branded PowerPoint (.pptx) from structured text.

Uses Cognizant_Template.pptx by default (coded_tools/assets/).
Falls back to a blank 16:9 presentation if the template is not found.

Slide syntax (in the content string):
  === COVER: Title | Subtitle ===    â†’ Cover slide (dark bg, white text)
  === DIVIDER: Section Title ===     â†’ Section divider slide (dark gradient)
  === CLOSING: Message ===           â†’ Closing slide (light bg)
  === Title ===                      â†’ Content slide (light bg, default)
  - bullet text                      â†’ Bullet point on the current slide
  * bullet text                      â†’ Bullet point on the current slide
  Plain text                         â†’ Body paragraph on the current slide

Modes:
  mode="write"  (default) â€” create a fresh presentation from the Cognizant template
  mode="append"           â€” open an existing .pptx and add slides at the end;
                            if the file does not exist, creates it fresh

Template layout indices (Cognizant_Template.pptx):
  0  Closing                         â€” light bg, dark text
  1  Divider 1 - gradient 1          â€” dark gradient bg, white text
  2  Basic title and content - light â€” light bg, dark text  â† default content
  3  Cover 1                         â€” dark bg, white text

Cognizant brand colours:
  #000048  dark navy   (body text on light slides)
  #2E308E  accent blue
  #7373D8  purple
  #06C7CC  teal
  #FFFFFF  white

HOCON class reference
â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "class": "coded_tools.common.write_pptx.WritePptx"

HOCON tool block (copy-paste into any network)
â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    {
        "name": "write_pptx",
        "class": "coded_tools.common.write_pptx.WritePptx",
        "function": {
            "description": "Create a Cognizant-branded PowerPoint (.pptx) from structured text. Use === Title === for content slides, === COVER: T | S === for cover, === DIVIDER: T === for section breaks, === CLOSING: T === for the closing slide. Use - for bullets.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path":          { "type": "string", "description": "Relative path for the output file inside the output directory (e.g. 'outputs/proposal.pptx')." },
                    "content":       { "type": "string", "description": "Slide content. Use === COVER: Title | Subtitle ===, === DIVIDER: Title ===, === CLOSING: Message ===, === Slide Title === for slides. Use - or * for bullets. Plain text becomes body paragraphs." },
                    "mode":          { "type": "string", "description": "'write' (default â€” create fresh from Cognizant template) or 'append' (add slides to existing file; creates fresh if file does not exist)." },
                    "template_path": { "type": "string", "description": "Optional path to a custom .pptx template. Overrides the default Cognizant template (write mode only)." },
                    "agent":         { "type": "string", "description": "Calling agent name (audit log)." }
                },
                "required": ["path", "content"]
            }
        }
    }

sly_data keys read
â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    output_dir     (preferred) â€” where relative output paths are resolved
    workspace_dir  (fallback)
    project_folder (fallback)  â€” bidmagic/dealcraft compat

Template
â”€â”€â”€â”€â”€â”€â”€â”€â”€
    Default: coded_tools/assets/Cognizant_Template.pptx
    Fallback: blank 16:9 presentation (if template file is missing)
    Ignored in mode="append" â€” the existing file's layouts are used
"""

from __future__ import annotations

import asyncio
import logging
import os
import re
from typing import Any

from neuro_san.interfaces.coded_tool import CodedTool

from coded_tools.common._base import log_call, resolve_output_path

logger = logging.getLogger(__name__)

# Template path â€” one level up from coded_tools/common/
_ASSETS_DIR    = os.path.join(os.path.dirname(__file__), "..", "assets")
_PPTX_TEMPLATE = os.path.abspath(os.path.join(_ASSETS_DIR, "Cognizant_Template.pptx"))

# Layout indices in Cognizant_Template.pptx
L_CLOSING = 0
L_DIVIDER = 1
L_CONTENT = 2
L_COVER   = 3

# Slide type markers
_RE_COVER   = re.compile(r"^===\s*COVER\s*:\s*(.+?)\s*===$",   re.IGNORECASE)
_RE_DIVIDER = re.compile(r"^===\s*DIVIDER\s*:\s*(.+?)\s*===$", re.IGNORECASE)
_RE_CLOSING = re.compile(r"^===\s*CLOSING\s*:\s*(.+?)\s*===$", re.IGNORECASE)
_RE_CONTENT = re.compile(r"^===\s*(.+?)\s*===$")


def _resolve_template(override: str | None) -> str | None:
    """Return absolute path to a valid template file, or None."""
    if override:
        t = override if os.path.isabs(override) else os.path.abspath(override)
        if os.path.exists(t):
            return t
        logger.warning("WritePptx: custom template '%s' not found â€” trying default", t)
    if os.path.exists(_PPTX_TEMPLATE):
        return _PPTX_TEMPLATE
    logger.warning("WritePptx: Cognizant_Template.pptx not found at '%s' â€” blank prs", _PPTX_TEMPLATE)
    return None


def _get_layout(prs, idx: int):
    """Return layout by index, or the first available layout as fallback."""
    layouts = prs.slide_layouts
    if idx < len(layouts):
        return layouts[idx]
    return layouts[0]


def _add_text_to_placeholder(slide, ph_idx: int, text: str, bold: bool = False):
    """Try to set text on a placeholder by index. Silent if not found."""
    try:
        from pptx.util import Pt  # type: ignore
        ph = slide.placeholders[ph_idx]
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        if bold:
            run.font.bold = True
    except (KeyError, IndexError):
        pass


def _add_text_box(slide, text: str, left, top, width, height, font_size=18, bold=False, color=None):
    """Add a plain text box to a slide."""
    from pptx.util import Pt       # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore

    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _build_slides(prs, content: str, has_template: bool):
    """Parse content and add slides to the presentation."""
    from pptx.util import Inches, Pt  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore

    C_DARK  = RGBColor(0x00, 0x00, 0x48)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # Slide geometry
    SW = prs.slide_width
    SH = prs.slide_height

    lines = content.splitlines()

    # Current slide state
    current_slide = None
    current_tf = None      # text frame for bullet/body content
    slide_count = 0

    def _new_slide(layout_idx: int, title: str, subtitle: str = ""):
        nonlocal current_slide, current_tf, slide_count
        layout = _get_layout(prs, layout_idx)
        current_slide = prs.slides.add_slide(layout)
        slide_count += 1
        current_tf = None

        # Set title placeholder (idx 0) if available
        try:
            ph = current_slide.placeholders[0]
            ph.text = title
        except (KeyError, IndexError):
            pass

        # Set subtitle/body placeholder (idx 1) for cover
        if subtitle:
            try:
                ph = current_slide.placeholders[1]
                ph.text = subtitle
            except (KeyError, IndexError):
                pass

        # Body placeholder (idx 1 for content slides) â†’ use for bullets
        if layout_idx == L_CONTENT:
            try:
                current_tf = current_slide.placeholders[1].text_frame
                current_tf.clear()
            except (KeyError, IndexError):
                current_tf = None

    def _add_bullet(text: str):
        """Add a bullet to the current slide's body placeholder."""
        nonlocal current_tf
        if current_slide is None:
            # No slide yet â€” create a default content slide
            _new_slide(L_CONTENT, "")

        if current_tf is None:
            # Try to get/create a body text frame
            try:
                current_tf = current_slide.placeholders[1].text_frame
                current_tf.clear()
            except (KeyError, IndexError):
                # Add a text box as fallback
                body = current_slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), SW - Inches(1), SH - Inches(2)
                )
                current_tf = body.text_frame
                current_tf.word_wrap = True

        # Add paragraph
        paras = current_tf.paragraphs
        if paras and paras[0].text == "":
            p = paras[0]
        else:
            p = current_tf.add_paragraph()
        run = p.add_run()
        run.text = text

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        m_cover   = _RE_COVER.match(stripped)
        m_divider = _RE_DIVIDER.match(stripped)
        m_closing = _RE_CLOSING.match(stripped)
        m_content = _RE_CONTENT.match(stripped)

        if m_cover:
            parts = m_cover.group(1).split("|", 1)
            title    = parts[0].strip()
            subtitle = parts[1].strip() if len(parts) > 1 else ""
            _new_slide(L_COVER if has_template else 0, title, subtitle)

        elif m_divider:
            _new_slide(L_DIVIDER if has_template else 0, m_divider.group(1))

        elif m_closing:
            _new_slide(L_CLOSING if has_template else 0, m_closing.group(1))

        elif m_content:
            _new_slide(L_CONTENT if has_template else 0, m_content.group(1))

        elif stripped.startswith("- ") or stripped.startswith("* "):
            _add_bullet(stripped[2:])

        else:
            _add_bullet(stripped)

    return slide_count


class WritePptx(CodedTool):
    """Create a Cognizant-branded PowerPoint (.pptx) from structured text."""

    def invoke(self, args: dict[str, Any], sly_data: dict[str, Any]) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return "Error: python-pptx is not installed. Run: pip install python-pptx"

        path_raw          = (args.get("path")          or "").strip()
        content           =  args.get("content")
        mode_raw          = (args.get("mode")          or "write").strip().lower()
        template_override = (args.get("template_path") or "").strip() or None
        agent             = (args.get("agent")         or "unknown-agent").strip()

        if not path_raw:
            return "Error: write_pptx requires 'path'."
        if content is None:
            return "Error: write_pptx requires 'content'."
        if mode_raw not in ("write", "append"):
            return f"Error: 'mode' must be 'write' or 'append' (got {mode_raw!r})."
        if not path_raw.lower().endswith(".pptx"):
            path_raw = path_raw + ".pptx"

        try:
            abs_path = resolve_output_path(path_raw, sly_data)
        except ValueError as exc:
            log_call(sly_data, tool="WritePptx", agent=agent, target=path_raw,
                     status="ERROR", detail=str(exc))
            return f"Error: {exc}"

        try:
            os.makedirs(os.path.dirname(abs_path), exist_ok=True)

            if mode_raw == "append" and os.path.exists(abs_path):
                # Open existing presentation and add slides at the end
                prs = Presentation(abs_path)
                has_template = True   # the file itself provides layouts
                mode_note = "append"
            else:
                # Create fresh from template (write mode, or append on missing file)
                template_file = _resolve_template(template_override)
                has_template  = template_file is not None
                prs = Presentation(template_file) if has_template else Presentation()

                # Remove all existing slides from the template
                if has_template:
                    xml_slides = prs.slides._sldIdLst
                    for _sld in list(prs.slides):
                        rId = xml_slides[0].get("r:id")
                        if rId:
                            prs.part.drop_rel(rId)
                        xml_slides.remove(xml_slides[0])

                mode_note = "write" if mode_raw == "write" else "append-as-create"

            slide_count = _build_slides(prs, content, has_template)
            prs.save(abs_path)
            file_size = os.path.getsize(abs_path)

        except Exception as exc:
            log_call(sly_data, tool="WritePptx", agent=agent, target=path_raw,
                     status="ERROR", detail=str(exc))
            logger.exception("WritePptx failed for %s", path_raw)
            return f"Error: failed to write '{path_raw}': {exc}"

        if mode_note == "append":
            template_note = "appended to existing file"
        elif has_template:
            template_note = "Cognizant template"
        else:
            template_note = "plain (template not found)"

        log_call(sly_data, tool="WritePptx", agent=agent, target=path_raw,
                 status="OK", detail=f"{mode_note}: {slide_count} slides, {file_size} bytes ({template_note})")
        logger.debug("WritePptx: %s (%s, %d slides, %d bytes, %s)",
                    path_raw, mode_note, slide_count, file_size, template_note)

        verb = "appended slides to" if mode_note == "append" else "wrote"
        return (
            f"OK: {verb} '{path_raw}' â€” "
            f"{slide_count} slides, {file_size} bytes ({template_note})."
        )

    async def async_invoke(self, args: dict[str, Any], sly_data: dict[str, Any]) -> str:
        return await asyncio.to_thread(self.invoke, args, sly_data)

