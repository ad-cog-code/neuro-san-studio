"""
coded_tools/common/read_pptx.py
════════════════════════════════
ReadPptx — extract text from a PowerPoint (.pptx) file, slide by slide.

Each slide returned as "=== Slide N: Title ===" with body text and notes.
Use start_slide / end_slide for large decks.
When truncated, the response tells the agent exactly what to call next.

Path resolution:
  • Relative paths → resolved against input_dir (sly_data)
  • Absolute paths → used as-is (Flask upload folders)

HOCON class reference
──────────────────────
    "class": "coded_tools.common.read_pptx.ReadPptx"

HOCON tool block (copy-paste into any network)
───────────────────────────────────────────────
    {
        "name": "read_pptx",
        "class": "coded_tools.common.read_pptx.ReadPptx",
        "function": {
            "description": "Extract text from a PowerPoint (.pptx) file. Returns per-slide content with titles and speaker notes. Use start_slide/end_slide for large decks — the response tells you the next start_slide when truncated.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path":          { "type": "string",  "description": "PPTX path. Relative to input_dir, or absolute." },
                    "start_slide":   { "type": "integer", "description": "1-based slide to start from (optional)." },
                    "end_slide":     { "type": "integer", "description": "1-based slide to stop at, inclusive (optional)." },
                    "include_notes": { "type": "boolean", "description": "Include speaker notes (default: true)." },
                    "agent":         { "type": "string",  "description": "Calling agent name (audit log)." }
                },
                "required": ["path"]
            }
        }
    }

sly_data keys read
───────────────────
    input_dir      (preferred) — where relative paths are resolved
    workspace_dir  (fallback)
    project_folder (fallback)  — bidmagic/dealcraft compat
"""

from __future__ import annotations

import asyncio
import logging
import os
from typing import Any

from neuro_san.interfaces.coded_tool import CodedTool
from coded_tools.common._base import log_call, resolve_input_path

logger = logging.getLogger(__name__)

MAX_RETURN_BYTES = 64_000


class ReadPptx(CodedTool):

    def invoke(self, args: dict[str, Any], sly_data: dict[str, Any]) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return "Error: python-pptx is not installed. Run: pip install python-pptx"

        path_raw      = (args.get("path")  or "").strip()
        agent         = (args.get("agent") or "unknown-agent").strip()
        start_slide   = args.get("start_slide")
        end_slide     = args.get("end_slide")
        include_notes = args.get("include_notes", True)

        if not path_raw:
            return "Error: read_pptx requires 'path'."

        abs_path = resolve_input_path(path_raw, sly_data)

        if not os.path.exists(abs_path):
            log_call(sly_data, tool="ReadPptx", agent=agent, target=path_raw,
                     status="MISS", detail="file not found")
            return f"NOT_FOUND: '{path_raw}' does not exist."

        if not abs_path.lower().endswith(".pptx"):
            return f"Error: '{path_raw}' does not appear to be a .pptx file."

        try:
            prs = Presentation(abs_path)
            total_slides = len(prs.slides)

            s = max(0, int(start_slide) - 1) if start_slide is not None else 0
            e = min(int(end_slide), total_slides) if end_slide is not None else total_slides
            slides = list(prs.slides)[s:e]

            parts = []
            for idx, slide in enumerate(slides, start=s + 1):
                title_text = ""
                body_lines = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    shape_text = shape.text_frame.text.strip()
                    if not shape_text:
                        continue
                    is_title = (
                        shape.shape_id == 1
                        or (
                            hasattr(shape, "placeholder_format")
                            and shape.placeholder_format is not None
                            and shape.placeholder_format.idx == 0
                        )
                    )
                    if is_title:
                        title_text = shape_text
                    else:
                        body_lines.append(shape_text)

                header = f"=== Slide {idx}"
                if title_text:
                    header += f": {title_text}"
                header += " ==="
                parts.append(header)

                if body_lines:
                    parts.extend(body_lines)

                if include_notes and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes: {notes}]")

                parts.append("")

        except Exception as exc:
            log_call(sly_data, tool="ReadPptx", agent=agent, target=path_raw,
                     status="ERROR", detail=str(exc))
            return f"Error: failed to read '{path_raw}': {exc}"

        content = "\n".join(parts)
        encoded = content.encode("utf-8")

        truncated    = False
        last_slide_in = e

        if len(encoded) > MAX_RETURN_BYTES:
            # Find how many slide blocks fit
            slide_texts = []
            current: list[str] = []
            for line in parts:
                if line.startswith("=== Slide ") and current:
                    slide_texts.append("\n".join(current))
                    current = [line]
                else:
                    current.append(line)
            if current:
                slide_texts.append("\n".join(current))

            cumulative = 0
            fits = 0
            for st in slide_texts:
                b = len(st.encode("utf-8")) + 1
                if cumulative + b > MAX_RETURN_BYTES:
                    break
                cumulative += b
                fits += 1

            content      = "\n".join(slide_texts[:fits])
            last_slide_in = s + fits
            truncated    = True

        sliced = start_slide is not None or end_slide is not None
        detail = f"{total_slides} slides total"
        if sliced:
            detail += f", requested slides {s + 1}–{e}"
        if truncated:
            next_s    = last_slide_in + 1
            remaining = e - last_slide_in
            detail   += f", truncated after slide {last_slide_in}"
            content  += (
                f"\n\n[TRUNCATED — {remaining} slide(s) not shown. "
                f"Call again with start_slide={next_s}"
                + (f", end_slide={e}" if end_slide is not None else "")
                + "]"
            )

        log_call(sly_data, tool="ReadPptx", agent=agent, target=path_raw,
                 status="OK", detail=detail)
        return content

    async def async_invoke(self, args: dict[str, Any], sly_data: dict[str, Any]) -> str:
        return await asyncio.to_thread(self.invoke, args, sly_data)
